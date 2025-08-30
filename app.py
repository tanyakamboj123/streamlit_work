"""
Text → PowerPoint Web App (Streamlit + python-pptx)

Features
- Paste long text/markdown
- Optional guidance ("investor pitch", "class lecture", etc.)
- Bring-your-own LLM key (OpenAI / Anthropic / Google Gemini)
- Upload .pptx/.potx template; app infers & reuses layouts, fonts, colors, and existing images
- Smart slide planning via LLM → JSON outline
- Builds a new .pptx using the uploaded template’s masters/layouts
- Reuses template images (no AI-generated images)
- Downloads the generated deck

Run locally
  pip install streamlit python-pptx requests pydantic
  streamlit run app.py
"""

import io
import json
import logging
import re
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

import requests
import streamlit as st
from pydantic import BaseModel, Field, ValidationError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

# -----------------------------
# App Config / Safety
# -----------------------------
logging.getLogger().setLevel(logging.CRITICAL)  # avoid logging sensitive info
st.set_page_config(page_title="Text → PowerPoint", page_icon="📑", layout="wide")

st.markdown(
    """
    <style>
    .small { font-size: 0.85rem; opacity: 0.8; }
    .ok { color: #0a8; }
    .warn { color: #a60; }
    .err { color: #c33; }
    </style>
    """,
    unsafe_allow_html=True,
)

# -----------------------------
# Models
# -----------------------------
class SlidePlan(BaseModel):
    title: str
    bullets: List[str] = Field(default_factory=list)
    layout_hint: Optional[str] = Field(
        default=None,
        description="e.g., 'Title and Content', 'Section Header', 'Title Only', 'Two Content'",
    )
    notes: Optional[str] = None

class DeckPlan(BaseModel):
    deck_title: Optional[str] = None
    slides: List[SlidePlan]

# -----------------------------
# LLM Callers (BYO key, never stored)
# -----------------------------

# Unified wrapper
class LLMProvider:
    OPENAI = "OpenAI"
    ANTHROPIC = "Anthropic"
    GEMINI = "Gemini"

@dataclass
class LLMRequest:
    provider: str
    api_key: str
    text: str
    guidance: str
    max_slides: int

LLM_SYSTEM_PROMPT = (
    "You are a world-class presentation architect. "
    "Given raw text or markdown and an optional guidance line, produce a JSON plan for a professional slide deck. "
    "Prioritize clarity, scannability, short bullets (6-10 words each), and logical sections. "
    "Do NOT include images. Do NOT exceed max_slides. Output STRICT JSON only."
)

LLM_USER_PROMPT_TMPL = (
    "INPUT_TEXT:\n{TEXT}\n\n"
    "GUIDANCE (may be empty): {GUIDANCE}\n\n"
    "Return JSON with shape:\n"
    "{\n  \"deck_title\": string | null,\n  \"slides\": [\n    {\n      \"title\": string,\n      \"bullets\": string[],\n      \"layout_hint\": string | null,\n      \"notes\": string | null\n    }\n  ]\n}\n\n"
    "Rules:\n"
    "- slides length <= {MAX_SLIDES}.\n"
    "- Use 'Section Header' for section dividers when helpful.\n"
    "- Keep bullets punchy; prefer 3-6 bullets per slide.\n"
    "- Use 'Two Content' only if the content clearly splits in two.\n"
    "- No markdown fences, no prose; JSON ONLY."
)

def _truncate(text: str, max_chars: int = 24000) -> Tuple[str, Optional[str]]:
    if len(text) <= max_chars:
        return text, None
    head = text[:max_chars]
    note = (
        f"Input truncated to {max_chars} characters to fit model context. "
        "Consider refining or using a model with larger context."
    )
    return head, note


def call_openai_chat(req: LLMRequest) -> str:
    url = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {req.api_key}", "Content-Type": "application/json"}
    prompt = LLM_USER_PROMPT_TMPL.format(TEXT=req.text, GUIDANCE=req.guidance, MAX_SLIDES=req.max_slides)
    payload = {
        "model": "gpt-4o-mini",
        "temperature": 0.2,
        "messages": [
            {"role": "system", "content": LLM_SYSTEM_PROMPT},
            {"role": "user", "content": prompt},
        ],
    }
    r = requests.post(url, headers=headers, json=payload, timeout=60)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"]


def call_anthropic(req: LLMRequest) -> str:
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        "x-api-key": req.api_key,
        "content-type": "application/json",
        "anthropic-version": "2023-06-01",
    }
    prompt = LLM_USER_PROMPT_TMPL.format(TEXT=req.text, GUIDANCE=req.guidance, MAX_SLIDES=req.max_slides)
    payload = {
        "model": "claude-3-5-sonnet-20240620",
        "max_tokens": 3000,
        "temperature": 0.2,
        "system": LLM_SYSTEM_PROMPT,
        "messages": [{"role": "user", "content": prompt}],
    }
    r = requests.post(url, headers=headers, json=payload, timeout=60)
    r.raise_for_status()
    return r.json()["content"][0]["text"]


def call_gemini(req: LLMRequest) -> str:
    # Gemini REST (Generative Language API)
    # https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-pro-latest:generateContent
    url = (
        "https://generativelanguage.googleapis.com/v1beta/models/"
        "gemini-1.5-pro-latest:generateContent?key=" + req.api_key
    )
    prompt = LLM_USER_PROMPT_TMPL.format(TEXT=req.text, GUIDANCE=req.guidance, MAX_SLIDES=req.max_slides)
    payload = {
        "contents": [
            {"parts": [{"text": LLM_SYSTEM_PROMPT + "\n\n" + prompt}]}
        ],
        "generationConfig": {"temperature": 0.2},
    }
    r = requests.post(url, json=payload, timeout=60)
    r.raise_for_status()
    # Some Gemini responses wrap with safety and candidates; extract text
    data = r.json()
    text = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")
    return text


def plan_slides(provider: str, api_key: str, text: str, guidance: str, max_slides: int) -> DeckPlan:
    text, trunc_note = _truncate(text)
    req = LLMRequest(provider=provider, api_key=api_key, text=text, guidance=guidance, max_slides=max_slides)
    if provider == LLMProvider.OPENAI:
        raw = call_openai_chat(req)
    elif provider == LLMProvider.ANTHROPIC:
        raw = call_anthropic(req)
    elif provider == LLMProvider.GEMINI:
        raw = call_gemini(req)
    else:
        raise ValueError("Unsupported provider")

    # Strip code fences if any
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?|```$", "", raw, flags=re.MULTILINE)
    try:
        obj = json.loads(raw)
        plan = DeckPlan(**obj)
    except Exception as e:
        raise RuntimeError(f"Failed to parse LLM JSON: {e}\nRaw: {raw[:800]}")

    # Optionally surface truncation note in a first slide
    if trunc_note:
        note_slide = SlidePlan(
            title="Note on Input Length",
            bullets=[trunc_note],
            layout_hint="Title and Content",
        )
        plan.slides = [note_slide] + plan.slides
        plan.slides = plan.slides[:max_slides]
    return plan

# -----------------------------
# PPT Helpers
# -----------------------------

def delete_all_slides(prs: Presentation) -> None:
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001 private API use is required
    slide_count = len(prs.slides)
    for _ in range(slide_count):
        sldIdLst.remove(sldIdLst[0])


def find_layout(prs: Presentation, hint: Optional[str]) -> Any:
    if not hint:
        return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    hint_l = hint.lower()
    # common names in Office themes
    for idx, layout in enumerate(prs.slide_layouts):
        name = (layout.name or "").lower()
        if hint_l in name:
            return layout
    # try fuzzy contains for well-known options
    mapping = {
        "title and content": ["title and content", "content"],
        "section header": ["section header", "section", "divider"],
        "title only": ["title only", "title"],
        "two content": ["two content", "comparison", "two column"],
    }
    for key, aliases in mapping.items():
        if hint_l == key or any(a in hint_l for a in aliases):
            for layout in prs.slide_layouts:
                nm = (layout.name or "").lower()
                if any(a in nm for a in aliases):
                    return layout
    # fallback
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


def collect_template_images(prs: Presentation) -> List[bytes]:
    seen = set()
    blobs: List[bytes] = []
    for s in prs.slides:
        for shp in s.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shp.image.blob
                    key = hash(blob)
                    if key not in seen:
                        seen.add(key)
                        blobs.append(blob)
                except Exception:
                    pass
    return blobs


def add_picture_decor(slide, blob: bytes) -> None:
    # decor: small image bottom-right at ~1.0" height
    stream = io.BytesIO(blob)
    slide_width = slide.part.slide_width
    slide_height = slide.part.slide_height
    pic = slide.shapes.add_picture(stream, left=slide_width - Inches(2.2), top=slide_height - Inches(1.8), height=Inches(1.0))
    # keep aspect ratio by specifying height only
    _ = pic  # noqa


def fill_placeholders(slide, title: str, bullets: List[str]) -> None:
    # prefer standard placeholder indices: 0 title, 1 body
    # title
    for shp in slide.shapes:
        if shp.is_placeholder and shp.placeholder_format.type == 1:  # TITLE
            shp.text_frame.clear()
            shp.text_frame.text = title
            break
    else:
        if slide.shapes.title:
            slide.shapes.title.text_frame.text = title
        else:
            tx = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(1))
            tx.text_frame.text = title
            tx.text_frame.paragraphs[0].font.size = Pt(32)

    # body
    body_frame = None
    for shp in slide.shapes:
        if shp.is_placeholder and shp.has_text_frame and shp.placeholder_format.idx in (1, 2):
            body_frame = shp.text_frame
            break
    if not body_frame:
        for shp in slide.shapes:
            if shp.has_text_frame and shp != slide.shapes.title:
                body_frame = shp.text_frame
                break
    if not body_frame:
        # last resort create a textbox
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(4.5))
        body_frame = tb.text_frame

    body_frame.clear()
    if bullets:
        body_frame.text = bullets[0]
        p0 = body_frame.paragraphs[0]
        p0.level = 0
        for b in bullets[1:]:
            p = body_frame.add_paragraph()
            p.text = b
            p.level = 0


def add_notes(slide, notes: Optional[str]):
    if not notes:
        return
    notes_slide = slide.notes_slide or slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes


def build_deck(template_bytes: bytes, plan: DeckPlan, reuse_images: bool = True) -> bytes:
    prs = Presentation(io.BytesIO(template_bytes))

    # Capture template images before clearing slides
    template_imgs = collect_template_images(prs) if reuse_images else []

    # Clear all slides to start clean while preserving masters/layouts/themes
    delete_all_slides(prs)

    # Optional title slide
    if plan.deck_title:
        layout = find_layout(prs, "Title")
        s = prs.slides.add_slide(layout)
        fill_placeholders(s, plan.deck_title, [])
        if template_imgs:
            add_picture_decor(s, template_imgs[0])

    # Content slides
    for sp in plan.slides:
        layout = find_layout(prs, sp.layout_hint or "Title and Content")
        s = prs.slides.add_slide(layout)
        fill_placeholders(s, sp.title, sp.bullets)
        add_notes(s, sp.notes)
        if template_imgs:
            # rotate through available images for visual consistency
            blob = template_imgs[(plan.slides.index(sp)) % len(template_imgs)]
            add_picture_decor(s, blob)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# -----------------------------
# UI
# -----------------------------

st.title("Your Text, Your Style – Auto‑Generate a Presentation")
st.caption("Paste text, pick style via your own PowerPoint template, and get a ready‑to‑present deck. No AI images.")

left, right = st.columns([1.1, 0.9])

with left:
    provider = st.radio("Choose LLM provider", [LLMProvider.OPENAI, LLMProvider.ANTHROPIC, LLMProvider.GEMINI], horizontal=True)
    api_key = st.text_input("Enter your API key (never stored)", type="password", help="Used only in-memory for this request; not logged or saved.")
    guidance = st.text_input("Optional guidance (e.g., 'turn into an investor pitch deck', 'client proposal', 'lecture slides')", "")
    max_slides = st.slider("Max slides", 4, 40, 14, help="Upper bound; the LLM will choose a reasonable number up to this.")

    raw_text = st.text_area(
        "Paste your text or markdown",
        height=280,
        placeholder="Paste long-form prose, notes, or markdown here…",
    )

with right:
    template_file = st.file_uploader("Upload a PowerPoint template or deck (.pptx or .potx)", type=["pptx", "potx"])  # noqa: RUF100
    st.markdown(
        "<span class='small'>Tip: pick a branded template to inherit fonts, colors, and layout. Existing images (logos, motifs) will be reused for visual consistency.</span>",
        unsafe_allow_html=True,
    )

    st.divider()
    st.markdown("**Privacy & Safety**")
    st.markdown(
        "<span class='small'>Your API key is never stored or logged. Files are processed in-memory and discarded after generation.</span>",
        unsafe_allow_html=True,
    )

st.divider()

btn = st.button("Generate Presentation", type="primary")

if btn:
    if not api_key:
        st.error("Please enter your API key for the selected provider.")
        st.stop()
    if not raw_text or len(raw_text.strip()) < 10:
        st.error("Please paste some source text (≥ 10 characters).")
        st.stop()
    if not template_file:
        st.error("Please upload a .pptx or .potx template/deck.")
        st.stop()

    with st.spinner("Planning slides…"):
        try:
            plan = plan_slides(provider, api_key, raw_text, guidance, max_slides)
        except (requests.HTTPError, RuntimeError, ValidationError) as e:
            st.error(f"LLM planning failed: {e}")
            st.stop()

    with st.spinner("Building deck from your template…"):
        try:
            deck_bytes = build_deck(template_file.read(), plan, reuse_images=True)
        except Exception as e:
            st.error(f"PowerPoint assembly failed: {e}")
            st.stop()

    st.success("Done! Your presentation is ready.")
    default_name = (plan.deck_title or "Generated Deck").strip() or "Generated Deck"
    default_name = re.sub(r"[^A-Za-z0-9 _-]", "", default_name)[:60]
    filename = f"{default_name}.pptx"

    st.download_button("⬇️ Download .pptx", data=deck_bytes, file_name=filename, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    st.markdown("---")
    with st.expander("Preview JSON plan (for debugging)"):
        st.code(plan.model_dump_json(indent=2), language="json")

st.markdown(
    """
**Notes & Limitations**
- The app inherits template styles by building slides using your uploaded deck’s masters/layouts. Exact 1:1 style replication depends on layout placeholder availability.
- Image reuse pulls any pictures found in your template and places small decorative motifs on slides. Adjust placement/size later in PowerPoint if needed.
- Extremely long inputs are truncated to stay within model context; refine guidance or raise max context with your chosen model.
- No images are generated via AI. Only images embedded in your template are reused.
"""
)
